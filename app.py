import streamlit as st
from pptx import Presentation
import fitz  
from langchain_openai import ChatOpenAI
from langchain_community.tools.tavily_search import TavilySearchResults
import json
import os
import re

st.set_page_config(page_title="AI Presentation Creator", page_icon="📊")

def extract_text_from_pptx(file):
    try:
        prs = Presentation(file)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        return "\n".join(text_content)
    except Exception as e:
        return f"Error reading PPTX: {e}"

def extract_text_from_pdf(file):
    try:
        doc = fitz.open(stream=file.read(), filetype="pdf")
        return "\n".join([page.get_text() for page in doc])
    except Exception as e:
        return f"Error reading PDF: {e}"

def create_pptx(slides_data):
    prs = Presentation()
    for slide_info in slides_data:
        # Use Title and Content layout (index 1)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        slide.shapes.title.text = slide_info.get("title", "Presentation Slide")
        
        # Set Content
        body_shape = slide.placeholders[1]
        body_shape.text = slide_info.get("content", "")
    
    output_path = "generated_presentation.pptx"
    prs.save(output_path)
    return output_path

# --- UI Layout ---
st.title("📊 AI Presentation Generator")

with st.sidebar:
    st.header("API Configuration")
    o_api = st.text_input("OpenAI API Key", type="password")
    t_api = st.text_input("Tavily API Key", type="password")
    num_slides = st.slider("Number of Slides", 5, 20, 10)
    st.divider()
    st.markdown("### How to use:")
    st.write("1. Enter Keys\n2. Enter Topic\n3. Upload PDF/PPTX (Optional)\n4. Click Generate")

topic = st.text_input("What is the presentation topic?", placeholder="e.g. The Future of Renewable Energy")
ref_file = st.file_uploader("Optional: Upload reference (PDF or PPTX)", type=["pdf", "pptx"])

if st.button("Generate Presentation"):
    if not o_api or not t_api or not topic:
        st.warning("⚠️ Please provide both API keys and a topic.")
    else:
        try:
            with st.spinner("Step 1: Gathering information..."):
                os.environ["TAVILY_API_KEY"] = t_api
                
                # 1. Handle reference file
                ref_content = ""
                if ref_file:
                    if ref_file.name.endswith("pptx"):
                        ref_content = extract_text_from_pptx(ref_file)
                    else:
                        ref_content = extract_text_from_pdf(ref_file)
                
                # 2. Search Web
                search = TavilySearchResults(max_results=3)
                web_data = search.invoke(topic)

            with st.spinner("Step 2: AI is writing your slides..."):
                # Use gpt-4o-mini for better rate limits
                llm = ChatOpenAI(model="gpt-4o-mini", api_key=o_api, temperature=0.7)
                
                # NOTE: We use double curly braces {{ }} to avoid the f-string 'Invalid format specifier' error
                prompt = f"""
                You are an expert presentation consultant.
                Create a presentation based on this topic: {topic}
                
                Reference Material from user: {ref_content[:1500]}
                Recent Web Research: {web_data}
                
                Create exactly {num_slides} slides. 
                Return the response ONLY as a JSON array of objects. 
                Each object must have "title" and "content" keys.
                
                Example format:
                [
                  {{
                    "title": "Slide Title Here",
                    "content": "Bullet point 1\\nBullet point 2\\nBullet point 3"
                  }}
                ]
                """
                
                res = llm.invoke(prompt)
                
                # Clean the response to ensure it's valid JSON
                clean_json_str = res.content.strip()
                if "```json" in clean_json_str:
                    clean_json_str = clean_json_str.split("```json")[1].split("```")[0].strip()
                elif "```" in clean_json_str:
                    clean_json_str = clean_json_str.split("```")[1].split("```")[0].strip()
                
                slides_json = json.loads(clean_json_str)
                
                # 3. Create PPTX
                pptx_path = create_pptx(slides_json)
                
                st.success(f"✅ Generated {len(slides_json)} slides successfully!")
                
                with open(pptx_path, "rb") as file:
                    st.download_button(
                        label="📥 Download PowerPoint File",
                        data=file,
                        file_name=f"{topic.replace(' ', '_')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )

        except json.JSONDecodeError:
            st.error("❌ The AI returned an invalid data format. Please try clicking Generate again.")
        except Exception as e:
            error_msg = str(e)
            if "insufficient_quota" in error_msg:
                st.error("❌ OpenAI Account Error: You have $0 balance. Please add credits at platform.openai.com.")
            elif "rate_limit" in error_msg.lower():
                st.error("❌ Rate Limit: You are making requests too fast or your account tier is too low.")
            else:
                st.error(f"⚠️ An error occurred: {error_msg}")
